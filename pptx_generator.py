"""
PPTX Generator — AI mazmuni asosida chiroyli PowerPoint fayl yasaydi
"""
import json
import random
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Color palettes (content-aware, rotated) ───────────
PALETTES = [
    {"primary": "1E2761", "secondary": "CADCFC", "accent": "FFFFFF", "name": "Midnight"},
    {"primary": "2C5F2D", "secondary": "97BC62", "accent": "FFFFFF", "name": "Forest"},
    {"primary": "065A82", "secondary": "8FD9E8", "accent": "FFFFFF", "name": "Ocean"},
    {"primary": "36454F", "secondary": "C8D0D6", "accent": "FFFFFF", "name": "Charcoal"},
    {"primary": "028090", "secondary": "8FE0D6", "accent": "FFFFFF", "name": "Teal"},
    {"primary": "990011", "secondary": "F4C7C7", "accent": "FFFFFF", "name": "Cherry"},
    {"primary": "6D2E46", "secondary": "D9B8C4", "accent": "FFFFFF", "name": "Berry"},
    {"primary": "B85042", "secondary": "F0D9C8", "accent": "FFFFFF", "name": "Terracotta"},
]

def hex_to_rgb(hex_color):
    return RGBColor.from_string(hex_color)

def set_background(slide, color_hex):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)

def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color="000000", align=PP_ALIGN.LEFT, font="Calibri", italic=False,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = hex_to_rgb(color)
    return box

def _no_shadow(shape):
    """Force-disable shadow: clear effectLst AND remove <p:style> (which re-applies theme shadow via effectRef)."""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is not None:
        for tag in ('a:effectLst', 'a:effectDag'):
            el = spPr.find(qn(tag))
            if el is not None:
                spPr.remove(el)
        from pptx.oxml import parse_xml
        empty_effect = parse_xml('<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
        spPr.append(empty_effect)
    # Remove <p:style> entirely — it carries effectRef/fillRef that reintroduce theme shadow
    style_el = sp.find(qn('p:style'))
    if style_el is not None:
        sp.remove(style_el)

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = hex_to_rgb(line_color)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    _no_shadow(shape)
    return shape

def add_icon_circle(slide, left, top, size, bg_color, icon_text, icon_color="FFFFFF", icon_size=20):
    circle = add_shape(slide, MSO_SHAPE.OVAL, left, top, size, size, fill_color=bg_color)
    tf = circle.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = icon_text
    run.font.size = Pt(icon_size)
    run.font.color.rgb = hex_to_rgb(icon_color)
    run.font.bold = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return circle

SLIDE_W = 13.333
SLIDE_H = 7.5

def build_title_slide(prs, palette, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, palette["primary"])
    # Decorative circle motif
    add_shape(slide, MSO_SHAPE.OVAL, SLIDE_W - 3.5, -2, 6, 6, fill_color=palette["secondary"])
    add_shape(slide, MSO_SHAPE.OVAL, -2, SLIDE_H - 2.5, 5, 5, fill_color=palette["secondary"])
    # dim overlay rectangle won't apply since pptx doesn't support opacity easily; keep clean
    add_text(slide, 1, SLIDE_H/2 - 1.4, SLIDE_W - 2, 1.8, title, size=44, bold=True,
              color=palette["accent"], align=PP_ALIGN.LEFT, font="Cambria", anchor=MSO_ANCHOR.BOTTOM)
    add_text(slide, 1, SLIDE_H/2 + 0.5, SLIDE_W - 2, 0.8, subtitle, size=18,
              color="E8E8E8", align=PP_ALIGN.LEFT, italic=True, anchor=MSO_ANCHOR.TOP)
    return slide

def build_content_slide(prs, palette, slide_num, total, heading, bullets, layout_variant=0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, "FFFFFF")

    # Header bar area (no stripe - use color block instead)
    add_text(slide, 0.6, 0.45, SLIDE_W - 1.2, 0.9, heading, size=30, bold=True,
              color=palette["primary"], font="Cambria")

    # slide number badge
    add_icon_circle(slide, SLIDE_W - 1.0, 0.45, 0.5, palette["primary"], str(slide_num), icon_size=14)

    content_top = 1.6
    variant = layout_variant % 3

    if variant == 0:
        # Icon + text rows, vertically centered
        n = len(bullets)
        row_h = min(1.15, max(0.8, (SLIDE_H - content_top - 0.6) / max(n, 1)))
        total_h = row_h * n
        avail_h = SLIDE_H - content_top - 0.6
        y = content_top + max(0, (avail_h - total_h) / 2)
        for i, b in enumerate(bullets):
            add_icon_circle(slide, 0.7, y + row_h/2 - 0.225, 0.45, palette["secondary"], "•", icon_color=palette["primary"], icon_size=20)
            add_text(slide, 1.4, y, SLIDE_W - 2.2, row_h, b, size=16, color="2B2B2B", anchor=MSO_ANCHOR.MIDDLE)
            y += row_h
    elif variant == 1:
        # Two-column cards, sized to content, vertically centered
        half = (len(bullets) + 1) // 2
        col1, col2 = bullets[:half], bullets[half:]
        card_w = (SLIDE_W - 1.8) / 2
        row_h = 0.85
        pad = 0.3
        max_lines = max(len(col1), len(col2))
        card_h = pad * 2 + row_h * max_lines
        avail_h = SLIDE_H - content_top - 0.6
        v_offset = content_top + max(0, (avail_h - card_h) / 2)
        for ci, col in enumerate([col1, col2]):
            x = 0.7 + ci * (card_w + 0.4)
            this_h = pad * 2 + row_h * len(col) if col else 0.5
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, v_offset, card_w, this_h, fill_color=palette["secondary"])
            ty = v_offset + pad
            for b in col:
                add_text(slide, x + 0.3, ty, card_w - 0.6, row_h, "• " + b, size=14, color=palette["primary"], anchor=MSO_ANCHOR.MIDDLE)
                ty += row_h
    else:
        # Numbered process flow, vertically centered
        n = len(bullets)
        row_h = 1.0
        total_h = row_h * n
        avail_h = SLIDE_H - content_top - 0.6
        y = content_top + max(0, (avail_h - total_h) / 2)
        for i, b in enumerate(bullets):
            num_circle_color = palette["primary"]
            add_icon_circle(slide, 0.8, y + row_h/2 - 0.275, 0.55, num_circle_color, str(i+1), icon_size=18)
            add_text(slide, 1.6, y, SLIDE_W - 2.4, row_h, b, size=16, color="2B2B2B", anchor=MSO_ANCHOR.MIDDLE)
            y += row_h

    # footer page indicator
    add_text(slide, SLIDE_W - 1.8, SLIDE_H - 0.5, 1.4, 0.35, f"{slide_num} / {total}",
              size=10, color="999999", align=PP_ALIGN.RIGHT)
    return slide

def build_conclusion_slide(prs, palette, title, points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, palette["primary"])
    add_shape(slide, MSO_SHAPE.OVAL, SLIDE_W - 3, SLIDE_H - 3, 5, 5, fill_color=palette["secondary"])
    add_text(slide, 0.9, 0.7, SLIDE_W - 1.8, 1.0, title, size=36, bold=True,
              color=palette["accent"], font="Cambria")
    y = 2.0
    for p in points:
        add_icon_circle(slide, 0.9, y + 0.03, 0.4, palette["secondary"], "✓", icon_color=palette["primary"], icon_size=16)
        add_text(slide, 1.5, y, SLIDE_W - 2.6, 0.7, p, size=16, color=palette["accent"], anchor=MSO_ANCHOR.MIDDLE)
        y += 0.75
    add_text(slide, 0.9, SLIDE_H - 0.9, SLIDE_W - 1.8, 0.5, "Diqqatingiz uchun rahmat!",
              size=14, italic=True, color=palette["secondary"])
    return slide

def create_presentation(topic: str, slides_data: list, output_path: str):
    """
    slides_data: list of dicts like:
    [
      {"heading": "...", "bullets": ["...", "..."]},
      ...
    ]
    First item treated as title info if it has 'subtitle' key.
    Last item rendered as conclusion if marked is_conclusion=True
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    palette = random.choice(PALETTES)

    total = len(slides_data)
    # Title slide
    first = slides_data[0]
    build_title_slide(prs, palette, topic, first.get("subtitle", "Taqdimot"))

    body_slides = slides_data[1:-1] if len(slides_data) > 2 else slides_data[1:]
    for i, sd in enumerate(body_slides):
        build_content_slide(prs, palette, i + 2, total, sd["heading"], sd["bullets"], layout_variant=i)

    if len(slides_data) > 2:
        last = slides_data[-1]
        build_conclusion_slide(prs, palette, last.get("heading", "Xulosa"), last.get("bullets", []))

    prs.save(output_path)
    return output_path
